import os
import requests
import uuid
import chardet
from io import BytesIO
from typing import Literal,Callable, Any, Optional,Union
from pydantic import BaseModel, Field, field_validator,PrivateAttr
from urllib.parse import urlparse
from pptx import Presentation

MAX_FILE_SIZE = 10 * 1024 * 1024

class File(BaseModel):
    """
    通用文件对象，支持自动类型推断和路径管理
    """
    url: str = Field(..., description="文件URL(http/https)或本地路径")
    file_type: Literal['image', 'video', 'audio', 'document', 'default'] = Field(
        default="default",
        description="文件类型"
    )
    _local_path: Optional[str] = PrivateAttr(default=None)

    def set_cache_path(self, path: str):
        """设置缓存路径"""
        self._local_path = path

    def get_cache_path(self) -> Optional[str]:
        """获取缓存路径（如果文件实际存在）"""
        return self._local_path

    @property
    def is_remote(self) -> bool:
        """判断是网络URL还是本地文件"""
        return self.url.startswith(('http://', 'https://'))

def infer_file_category(path_or_url: str) -> tuple[str, str]:
    """
    根据路径或URL后缀判断文件类型
    逻辑：
    1. 解析 URL 去除 query 参数 (?id=...)，提取 path
    2. 获取 path 最后一部分的文件名和后缀
    3. 查表判断，匹配不到则返回 'default'

    Return:
        - 分类:image, video, audio, document, default
        - 后缀:.pdf

    """

    # === 步骤 1 & 2: 提取纯净的后缀名 ===
    # urlparse 可以同时处理本地路径 (会被视为 path) 和 网络 URL
    parsed = urlparse(path_or_url)
    path = parsed.path  # 提取路径部分，忽略 http://... 和 ?query=...

    # 获取文件名 (例如 /a/b/test.jpg -> test.jpg)
    filename = os.path.basename(path)

    # 分离后缀 (test.jpg -> .jpg)
    _, ext_with_dot = os.path.splitext(filename)

    # 如果没有后缀，直接兜底
    if not ext_with_dot:
        return 'default', ""

    # 去除点并转小写 (例如 .JPG -> jpg)
    ext = ext_with_dot.lstrip('.').lower()

    # === 步骤 3: 查表匹配 ===
    # 定义常见映射表
    TYPE_MAPPING = {
        'image': {
            'apng', 'avif', 'bmp', 'gif', 'heic', 'ico', 'jpg', 'jpeg', 'png', 'svg', 'tiff', 'webp'
        },
        'video': {
            'mp4', 'avi', 'mov', 'mkv', 'flv', 'wmv', 'webm', 'm4v', '3gp'
        },
        'audio': {
            'mp3', 'wav', 'flac', 'aac', 'ogg', 'wma', 'm4a'
        },
        'document': {
            'pdf', 'doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx',
            'txt', 'md', 'csv', 'json', 'xml', 'html', 'htm'
        },
    }

    for category, extensions in TYPE_MAPPING.items():
        if ext in extensions:
            return category, ext_with_dot

    return 'default', ext_with_dot

class FileOps:
    DOWNLOAD_DIR = "/tmp"

    @staticmethod
    def read_content(file_obj:File, max_length=10000) -> str:
        return ""

    @staticmethod
    def get_local_path(file_obj:File) -> str:
        return file_obj.url

    @staticmethod
    def _get_bytes_stream(file_obj:File) -> tuple[bytes, str]:
        """
        获取文件内容和后缀, 5MB大小限制检查, 超出抛异常
        """
        _, ext = infer_file_category(file_obj.url)

        if file_obj.is_remote:
            try:
                # stream=True: 此时只下载 Headers，连接保持打开，还没下载 Body
                with requests.get(file_obj.url, stream=True, timeout=60) as resp:
                    resp.raise_for_status()

                    content_length = resp.headers.get('Content-Length')
                    if content_length and int(content_length) > MAX_FILE_SIZE:
                        raise Exception(
                            f"文件大小 ({int(content_length)} bytes) 超过限制 5MB，已终止下载。"
                        )

                    # 场景：Header 缺失 Content-Length 或服务器 Header 欺骗
                    downloaded_content = BytesIO()
                    current_size = 0

                    # 分块读取，每块 8KB
                    for chunk in resp.iter_content(chunk_size=8192):
                        if chunk:
                            current_size += len(chunk)
                            if current_size > MAX_FILE_SIZE:
                                raise Exception(f"检测到文件超过 5MB，已中断。")
                            downloaded_content.write(chunk)

                    # 获取完整 bytes
                    return downloaded_content.getvalue(), ext

            except requests.RequestException as e:
                raise RuntimeError(f"网络请求失败: {e}")

        else:
            if not os.path.exists(file_obj.url):
                raise FileNotFoundError(f"本地文件不存在: {file_obj.url}")

            '''
            file_size = os.path.getsize(file_obj.url)
            if file_size > MAX_FILE_SIZE:
                 raise Exception(f"本地文件大小 ({file_size} bytes) 超过限制 5MB")
            '''

            with open(file_obj.url, 'rb') as f:
                return f.read(), ext

    @staticmethod
    def save_to_local(file_obj: File, filename: str) -> str:
        """
        将当前文件对象的内容保存到本地路径, 返回本地路径
        如果是本地路径，直接返回
        """
        if not file_obj.is_remote:
            if os.path.exists(file_obj.url):
                return file_obj.url

            raise FileNotFoundError(f"Local file not found: {file_obj.url}")

        try:
            os.makedirs(FileOps.DOWNLOAD_DIR, exist_ok=True)

            # 简单的文件名生成策略 (真实场景建议用 url hash 避免重复下载)
            # ext = os.path.splitext(file_obj.url.split('?')[0])[1] or ".tmp"
            # filename = f"{uuid.uuid4().hex}{ext}"
            local_path = os.path.join(FileOps.DOWNLOAD_DIR, filename)

            headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3'}
            with requests.get(file_obj.url, headers=headers, stream=True, timeout=120) as r:
                r.raise_for_status()
                with open(local_path, 'wb') as f:
                    for chunk in r.iter_content(chunk_size=8192):
                        f.write(chunk)

            return local_path
        except Exception as e:
            raise RuntimeError(f"Download failed for {file_obj.url}: {str(e)}")

    @staticmethod
    def read_bytes(file_obj:File) -> bytes:
        """
        获取文件的原始二进制数据
        场景：上传到OSS、保存到本地、传给图像处理库
        """
        content, _ = FileOps._get_bytes_stream(file_obj)
        return content

    @staticmethod
    def extract_text(file_obj: File) -> str:
        """
        提取文本内容
        场景：RAG、HTML解析、文档分析
        """
        try:
            content, ext = FileOps._get_bytes_stream(file_obj)

            if ext in ['.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx']:
                return FileOps._parse_document_bytes(file_obj, content, ext)

            # 默认直接读
            charset = chardet.detect(content)
            if 'encoding' in charset:
                return content.decode(charset['encoding'])
            else:
                return content.decode('utf-8')

        except Exception as e:
            return f"[FileOps Error] Failed to read content: {str(e)}"

    @staticmethod
    def _parse_document_bytes(file_obj: File, content: bytes, ext:str) -> str:
        stream = BytesIO(content)
        text_result = ""

        try:
            if ext == '.pdf':
                import pypdf
                reader = pypdf.PdfReader(stream)
                for page in reader.pages:
                    text_result += page.extract_text() + "\n"
            elif ext in ['.docx', '.doc']:
                text_result = read_docx(stream)
            elif ext in ['.xlsx', '.xls', '.csv']:
                import pandas as pd
                if ext == '.csv':
                    df = pd.read_csv(stream)
                else:
                    df = pd.read_excel(stream)
                text_result = df.to_string()
            elif ext in ['.ppt', '.pptx']:
                text_result = read_ppt(stream)
            else:
                text_result = f"[暂不支持解析该文档格式: {ext}]"
        except ImportError as e:
            text_result = f"[解析库缺失] {e}"
        except Exception as e:
            text_result = f"[解析失败] {e}"

        return text_result

def read_docx(cont_stream) -> str:
    """
    使用docx2python按顺序读取内容
    """
    from docx2python import docx2python
    doc_result = docx2python(cont_stream)

    # 获取文档结构
    all_parts = []

    # docx2python以嵌套列表形式返回内容
    # 遍历文档主体
    for section in doc_result.body:
        if isinstance(section, list):
            for item in section:
                if isinstance(item, list):
                    # 可能是表格或多级内容
                    for sub_item in item:
                        if isinstance(sub_item, str) and sub_item.strip():
                            all_parts.append(sub_item.strip())
                        elif isinstance(sub_item, list):
                            # 表格行
                            row_text = "\n".join([str(cell).strip() for cell in sub_item if str(cell).strip()])
                            if row_text:
                                all_parts.append(row_text)
                elif isinstance(item, str) and item.strip():
                    all_parts.append(item.strip())

    # 关闭文档
    doc_result.close()

    return "\n\n".join(all_parts)

def read_ppt(file_input: Union[str, bytes, BytesIO]) -> str:
    if not Presentation:
        return "[Error] 未安装 python-pptx 库，无法解析 PPT 文件"

    # 1. 统一转换为文件流对象 (BytesIO)
    if isinstance(file_input, str):
        with open(file_input, 'rb') as f:
            ppt_stream = BytesIO(f.read())
    elif isinstance(file_input, bytes):
        ppt_stream = BytesIO(file_input)
    else:
        ppt_stream = file_input

    try:
        prs = Presentation(ppt_stream)
        full_text = []

        for i, slide in enumerate(prs.slides):
            page_content = []
            page_content.append(f"=== 第 {i+1} 页 ===")

            # shape.text_frame 包含了形状内的文本段落
            for shape in slide.shapes:
                # 提取普通文本框
                if hasattr(shape, "text") and shape.text.strip():
                    page_content.append(shape.text.strip())

                # B. 提取表格内容 (普通 shape.text 无法获取表格内的字)
                if shape.has_table:
                    table_texts = []
                    for row in shape.table.rows:
                        row_cells = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame.text.strip()]
                        if row_cells:
                            table_texts.append(" | ".join(row_cells))
                    if table_texts:
                        page_content.append("[表格]\n" + "\n".join(table_texts))

            # 很多重要信息藏在备注里
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    page_content.append(f"[备注]: {notes.strip()}")

            full_text.append("\n".join(page_content))

        return "\n\n".join(full_text)

    except Exception as e:
        return f"[PPT解析失败] {str(e)}"
